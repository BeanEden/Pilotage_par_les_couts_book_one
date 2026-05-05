"""Script extraction PDF + PPTX — a lancer manuellement."""
import sys

try:
    from pptx import Presentation
    pptx_ok = True
except ImportError:
    pptx_ok = False
    print("python-pptx non installe. Lancez: pip install python-pptx")

try:
    from PyPDF2 import PdfReader
    pdf_ok = True
except ImportError:
    pdf_ok = False
    print("PyPDF2 non installe. Lancez: pip install PyPDF2")

if pdf_ok:
    print("=" * 60)
    print("CONTENU DU PDF : Rapport_DF_BookOne.pdf")
    print("=" * 60)
    reader = PdfReader("Rapport_DF_BookOne.pdf")
    for i, page in enumerate(reader.pages):
        print(f"\n--- Page {i+1} ---")
        print(page.extract_text())

if pptx_ok:
    print("\n" + "=" * 60)
    print("CONTENU DU PPTX : BookOne_MindMap_Slide.pptx")
    print("=" * 60)
    prs = Presentation("BookOne_MindMap_Slide.pptx")
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                print(shape.text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    print(" | ".join(cell.text for cell in row.cells))
