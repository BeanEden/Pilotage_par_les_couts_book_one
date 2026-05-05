"""Generateur PowerPoint — Presentation Book One MVP."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy, os

# Couleurs
DARK  = RGBColor(0x1E, 0x27, 0x61)
BLUE  = RGBColor(0x37, 0x8A, 0xDD)
GREEN = RGBColor(0x1D, 0x9E, 0x75)
RED   = RGBColor(0xE2, 0x4B, 0x4A)
ORANGE= RGBColor(0xBA, 0x75, 0x17)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0x5F, 0x5E, 0x5A)
LGRAY = RGBColor(0xF0, 0xF2, 0xF5)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_box(slide, left, top, w, h, text, font_size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT, bg=None):
    if bg:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.fill.background()
        tf = shape.text_frame
    else:
        txBox = slide.shapes.add_textbox(left, top, w, h)
        tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf

def add_multi(slide, left, top, w, h, lines, bg=None):
    """lines = [(text, size, bold, color), ...]"""
    if bg:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.fill.background()
        tf = shape.text_frame
    else:
        txBox = slide.shapes.add_textbox(left, top, w, h)
        tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, sz, bld, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.bold = bld
        p.font.color.rgb = col
        p.space_after = Pt(4)
    return tf

# ============================================================
# SLIDE 1 — Page de garde
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_box(sl, Inches(1), Inches(0.8), Inches(11), Inches(1), "BOOK ONE", 48, True, WHITE, PP_ALIGN.CENTER)
add_box(sl, Inches(1), Inches(1.8), Inches(11), Inches(0.8), "Tous les livres a 1 EUR — MVP Mobile", 28, False, BLUE, PP_ALIGN.CENTER)
add_box(sl, Inches(2), Inches(3.2), Inches(9), Inches(0.6), "Pilotage par les couts — Rapport de presentation", 20, False, RGBColor(0xCC,0xCC,0xCC), PP_ALIGN.CENTER)

# KPIs en bas
for i, (lab, val) in enumerate([("Duree", "8 semaines"), ("Budget", "130 000 EUR"), ("Equipe", "5 profils"), ("Statut", "MVP Livre")]):
    x = Inches(2 + i * 2.5)
    add_multi(sl, x, Inches(4.5), Inches(2), Inches(1.2),
              [(lab, 12, False, RGBColor(0xAA,0xAA,0xAA)), (val, 18, True, WHITE)],
              bg=RGBColor(0x2A, 0x33, 0x70))

add_box(sl, Inches(1), Inches(6.5), Inches(11), Inches(0.5), "Janvier — Fevrier 2025  |  Sup de Vinci M2", 14, False, GRAY, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 — Perimetre & Hypotheses
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_box(sl, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), "Perimetre & Hypotheses", 32, True, WHITE)

# MVP
add_multi(sl, Inches(0.5), Inches(1.3), Inches(3.8), Inches(2.5),
    [("MVP (9 fonctionnalites)", 16, True, WHITE),
     ("", 6, False, WHITE),
     ("Connexion / Inscription", 12, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Catalogue geolocise", 12, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Favoris & Wishes", 12, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Depot & Scan ISBN (3 livres)", 12, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Achat a 1 EUR", 12, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Messagerie RDV", 12, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Notifications vente/souhait", 12, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Mode invite", 12, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Remise en main propre", 12, False, RGBColor(0xBB,0xBB,0xBB)),
    ], bg=RGBColor(0x2A,0x33,0x70))

# Bonus
add_multi(sl, Inches(4.6), Inches(1.3), Inches(3.8), Inches(2.5),
    [("Bonus (4 fonctionnalites)", 16, True, ORANGE),
     ("", 6, False, WHITE),
     ("Mode invite etendu", 12, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Geolocalisation avancee", 12, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Integration Stripe", 12, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Remise en main propre", 12, False, RGBColor(0xBB,0xBB,0xBB)),
    ], bg=RGBColor(0x2A,0x33,0x70))

# Hypotheses
add_multi(sl, Inches(8.7), Inches(1.3), Inches(4.1), Inches(2.5),
    [("Hypotheses cles", 16, True, GREEN),
     ("", 6, False, WHITE),
     ("TJM/j x j/h = cout tache", 12, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Reserve imprevus 15%", 12, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Complexite : Faible / Moyenne / Haute", 12, False, RGBColor(0xBB,0xBB,0xBB)),
     ("5 jours ouvres par semaine", 12, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Charges patronales CDI : 42-45%", 12, False, RGBColor(0xBB,0xBB,0xBB)),
    ], bg=RGBColor(0x2A,0x33,0x70))

# Timeline
add_multi(sl, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.8),
    [("Debut 6 janvier 2025  —  Livraison 28 fevrier 2025  —  8 semaines", 16, True, BLUE)],
    bg=RGBColor(0x2A,0x33,0x70))

# Phases resume
phases = [("P1 Cadrage", "S1-S2", "5 200 EUR"), ("P2 Archi", "S1-S3", "12 800 EUR"),
          ("P3 Dev MVP", "S3-S6", "32 400 EUR"), ("P4 Bonus", "S5-S7", "14 500 EUR"),
          ("P5 Tests", "S6-S8", "9 800 EUR"), ("P6 Prod", "S7-S8", "5 700 EUR")]
for i, (nom, sem, cout) in enumerate(phases):
    x = Inches(0.5 + i * 2.1)
    add_multi(sl, x, Inches(5.3), Inches(1.9), Inches(1.8),
        [(nom, 13, True, WHITE), (sem, 11, False, BLUE), (cout, 14, True, GREEN)],
        bg=RGBColor(0x2A,0x33,0x70))

# ============================================================
# SLIDE 3 — Structure des couts
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_box(sl, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), "Structure des couts", 32, True, WHITE)

postes = [
    ("RH salaries CDI", "Chef projet, UX, Dev mob, Dev back, QA", "60 800 EUR", GREEN),
    ("Freelance", "Paiement, Securite, Geoloc, Archi", "12 600 EUR", BLUE),
    ("Alternant + Stagiaire", "80 EUR/j + 50 EUR/j (charges incluses)", "2 350 EUR", ORANGE),
    ("Technique & Outils", "Cloud, BDD, API, Stripe, Notifs, Monitoring", "5 050 EUR", RGBColor(0x7F,0x77,0xDD)),
    ("Couts caches", "Reunions, Corrections, Documentation", "2 250 EUR", GRAY),
    ("Mise en prod & stores", "Deploy + App Store + Play Store", "3 300 EUR", RED),
]
for i, (titre, detail, cout, col) in enumerate(postes):
    y = Inches(1.3 + i * 0.85)
    add_multi(sl, Inches(0.5), y, Inches(8), Inches(0.75),
        [(f"{titre}  —  {cout}", 15, True, col), (detail, 11, False, RGBColor(0xAA,0xAA,0xAA))],
        bg=RGBColor(0x2A,0x33,0x70))
    add_box(sl, Inches(8.8), y, Inches(2), Inches(0.75), cout, 18, True, col, PP_ALIGN.CENTER, bg=RGBColor(0x2A,0x33,0x70))

# Total
add_multi(sl, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7),
    [("Reserve imprevus 15% : 12 400 EUR   |   BUDGET TOTAL : 98 350 EUR", 18, True, WHITE)],
    bg=RGBColor(0x1D,0x9E,0x75))

# ============================================================
# SLIDE 4 — Mind Map (replique du PPTX existant)
# ============================================================
# On insere la slide existante directement
try:
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    src_prs = Presentation("BookOne_MindMap_Slide.pptx")
    src_slide = src_prs.slides[0]
    # Copie manuelle des shapes
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    # Copier le background
    src_bg = src_slide.background
    if src_bg.fill.type is not None:
        bg = sl.background
        fill = bg.fill
        fill.solid()
        try:
            fill.fore_color.rgb = src_bg.fill.fore_color.rgb
        except:
            fill.fore_color.rgb = DARK
    else:
        add_bg(sl)
    
    # Copier chaque shape
    for shape in src_slide.shapes:
        el = copy.deepcopy(shape._element)
        sl.shapes._spTree.append(el)
    print("Mind Map slide copiee avec succes.")
except Exception as e:
    print(f"Impossible de copier la slide MindMap: {e}")
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(sl)
    add_box(sl, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), "BOOK ONE — Budget Previsionnel MVP (Mind Map)", 32, True, WHITE)
    add_box(sl, Inches(1), Inches(2), Inches(11), Inches(4),
            "Voir le fichier BookOne_MindMap_Slide.pptx pour la version complete de la Mind Map.\n\n"
            "Budget total : 98 350 EUR  |  9 profils  |  168 j/h  |  3 mois\n\n"
            "RH : 73 400 EUR  |  Technique : 5 050 EUR  |  Couts caches : 2 250 EUR\n"
            "Mise en prod : 3 300 EUR  |  Reserve 15% : 12 400 EUR",
            16, False, RGBColor(0xBB,0xBB,0xBB))

# ============================================================
# SLIDE 5 — Equipe projet
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_box(sl, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), "Equipe & Ressources Humaines", 32, True, WHITE)

equipe = [
    ("PM", "Chef de projet", "CDI Cadre", "58 000 EUR", "7F77DD"),
    ("BE", "Dev Back-End Junior", "CDI ETAM", "38 000 EUR", "1D9E75"),
    ("MOB", "Dev Mobile Flutter", "CDI ETAM", "55 000 EUR", "378ADD"),
    ("QA", "QA Engineer", "CDI ETAM", "45 000 EUR", "D4537E"),
    ("FRL", "Freelance Paiement/Secu", "Freelance", "130 000 EUR", "E24B4A"),
]
for i, (rid, label, typ, sal, col) in enumerate(equipe):
    y = Inches(1.3 + i * 1.0)
    c = RGBColor(int(col[:2],16), int(col[2:4],16), int(col[4:],16))
    add_multi(sl, Inches(0.5), y, Inches(5.5), Inches(0.85),
        [(f"{rid}  —  {label}", 15, True, c),
         (f"{typ}  |  Salaire brut annuel : {sal}", 11, False, RGBColor(0xAA,0xAA,0xAA))],
        bg=RGBColor(0x2A,0x33,0x70))

add_multi(sl, Inches(7), Inches(1.3), Inches(5.8), Inches(2),
    [("Calcul du TJM", 16, True, BLUE),
     ("TJM = (Brut annuel x (1 + charges)) / Jours ouvres", 12, False, RGBColor(0xBB,0xBB,0xBB)),
     ("CDI Cadre : charges 45%  |  218 j/an", 12, False, RGBColor(0xBB,0xBB,0xBB)),
     ("CDI ETAM : charges 42%  |  218 j/an", 12, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Freelance : 0% charges  |  200 j/an", 12, False, RGBColor(0xBB,0xBB,0xBB)),
    ], bg=RGBColor(0x2A,0x33,0x70))

add_multi(sl, Inches(7), Inches(3.8), Inches(5.8), Inches(1.5),
    [("Levier d'optimisation", 16, True, ORANGE),
     ("Remplacement Dev Back-End Senior par Junior", 13, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Economie estimee : 6 500 EUR sur 8 semaines", 13, True, GREEN),
    ], bg=RGBColor(0x2A,0x33,0x70))

# ============================================================
# SLIDE 6 — Planning & Risques
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_box(sl, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), "Planning & Gestion des risques", 32, True, WHITE)

phases2 = [
    ("P1 — Cadrage & Architecture", "S1 - S2", "Backlog, stack technique, architecture cloud", "5 200 EUR"),
    ("P2 — UX/UI & Architecture", "S1 - S3", "Maquettes HD, CI/CD, modele de donnees", "12 800 EUR"),
    ("P3 — Developpement MVP", "S3 - S6", "API Back-End + App Mobile Flutter + Messagerie", "32 400 EUR"),
    ("P4 — Fonctionnalites Bonus", "S5 - S7", "Stripe, Geolocalisation, Mode invite, Remise", "14 500 EUR"),
    ("P5 — Tests & Recette", "S6 - S8", "QA fonctionnel + securite paiement + UAT", "9 800 EUR"),
    ("P6 — Mise en production", "S7 - S8", "Deploiement + App Store + Play Store + Monitoring", "5 700 EUR"),
]
for i, (nom, sem, desc, cout) in enumerate(phases2):
    y = Inches(1.2 + i * 0.8)
    add_multi(sl, Inches(0.5), y, Inches(8.5), Inches(0.7),
        [(f"{nom}  ({sem})", 14, True, BLUE), (desc, 11, False, RGBColor(0xAA,0xAA,0xAA))],
        bg=RGBColor(0x2A,0x33,0x70))
    add_box(sl, Inches(9.3), y, Inches(2), Inches(0.7), cout, 16, True, GREEN, PP_ALIGN.CENTER, bg=RGBColor(0x2A,0x33,0x70))

# Risques
add_box(sl, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.5), "Risques identifies", 18, True, RED)
risques = [
    "API ISBN : limite de taux depassee (Moyen — Resolu)",
    "Stripe : delai validation compte pro (Eleve — En cours)",
    "Derive charge Dev Mobile S7 (Moyen — En cours)",
    "DPO : recommandation stockage images (Faible — Resolu)",
]
add_multi(sl, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.8),
    [(r, 11, False, RGBColor(0xCC,0xCC,0xCC)) for r in risques],
    bg=RGBColor(0x2A,0x33,0x70))

# ============================================================
# SLIDE 7 — KPIs Direction Financiere
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_box(sl, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), "Dashboard Direction Financiere — KPIs", 32, True, WHITE)

kpis = [("Jalon", "S8"), ("Avancement", "100%"), ("Cout planifie", "31 590 EUR"),
        ("Cout reel", "33 422 EUR"), ("Ecart", "+5.8%"), ("Budget valide", "130 000 EUR")]
for i, (lab, val) in enumerate(kpis):
    x = Inches(0.5 + (i % 3) * 4.2)
    y = Inches(1.3 + (i // 3) * 1.5)
    col = GREEN if "100" in val or val.startswith("-") else RED if val.startswith("+") else BLUE
    add_multi(sl, x, y, Inches(3.8), Inches(1.2),
        [(lab, 13, False, RGBColor(0xAA,0xAA,0xAA)), (val, 24, True, col)],
        bg=RGBColor(0x2A,0x33,0x70))

# Jalons
add_box(sl, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.5), "Suivi des jalons", 18, True, WHITE)
jalons = [
    ("J1 — S4 : Maquettes CEO", "Prevu 18 000 EUR  |  Reel 17 600 EUR", "-400 EUR", GREEN),
    ("J2 — S8 : Livraison MVP", "Prevu 31 590 EUR  |  Reel 33 422 EUR", "+1 832 EUR", ORANGE),
]
for i, (nom, detail, ecart, col) in enumerate(jalons):
    y = Inches(5.2 + i * 1.0)
    add_multi(sl, Inches(0.5), y, Inches(8.5), Inches(0.85),
        [(nom, 15, True, col), (detail, 12, False, RGBColor(0xBB,0xBB,0xBB))],
        bg=RGBColor(0x2A,0x33,0x70))
    add_box(sl, Inches(9.3), y, Inches(3.5), Inches(0.85), ecart, 20, True, col, PP_ALIGN.CENTER, bg=RGBColor(0x2A,0x33,0x70))

# ============================================================
# SLIDE 8 — Synthese & Conclusion
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_box(sl, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), "Synthese & Conclusion", 32, True, WHITE)

add_multi(sl, Inches(0.5), Inches(1.3), Inches(6), Inches(3),
    [("Points cles du projet", 18, True, GREEN),
     ("", 6, False, WHITE),
     ("Budget maitrise : ecart < 6% sur le cout reel", 13, False, RGBColor(0xBB,0xBB,0xBB)),
     ("MVP livre en 8 semaines (jan-fev 2025)", 13, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Reserve de 15% preservee en grande partie", 13, False, RGBColor(0xBB,0xBB,0xBB)),
     ("5 profils optimises (levier Junior)", 13, False, RGBColor(0xBB,0xBB,0xBB)),
     ("10 KPIs de pilotage operationnels", 13, False, RGBColor(0xBB,0xBB,0xBB)),
    ], bg=RGBColor(0x2A,0x33,0x70))

add_multi(sl, Inches(6.8), Inches(1.3), Inches(6), Inches(3),
    [("Prochaines etapes", 18, True, BLUE),
     ("", 6, False, WHITE),
     ("Lancement beta test utilisateurs", 13, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Publication App Store & Play Store", 13, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Acquisition premiers utilisateurs", 13, False, RGBColor(0xBB,0xBB,0xBB)),
     ("Breakeven estime a M9-M10", 13, False, RGBColor(0xBB,0xBB,0xBB)),
     ("ROI positif a 24 mois", 13, False, RGBColor(0xBB,0xBB,0xBB)),
    ], bg=RGBColor(0x2A,0x33,0x70))

# Chiffres finaux
final = [("73 400 EUR", "Couts RH totaux"), ("5 050 EUR", "Technique & outils"),
         ("12 400 EUR", "Reserve 15%"), ("98 350 EUR", "Budget total")]
for i, (val, lab) in enumerate(final):
    x = Inches(0.5 + i * 3.2)
    add_multi(sl, x, Inches(5), Inches(2.8), Inches(1.5),
        [(val, 22, True, GREEN if i < 3 else WHITE), (lab, 12, False, RGBColor(0xAA,0xAA,0xAA))],
        bg=RGBColor(0x2A,0x33,0x70))

add_box(sl, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
        "Book One — Pilotage par les couts  |  Sup de Vinci M2  |  2025", 14, False, GRAY, PP_ALIGN.CENTER)

# ============================================================
# SAUVEGARDE
# ============================================================
output = "BookOne_Presentation.pptx"
prs.save(output)
print(f"Presentation generee : {output}")
print(f"  {len(prs.slides)} slides creees.")
